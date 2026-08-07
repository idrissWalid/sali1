import base64
import io
from unittest.mock import patch

import pytest
from PIL import Image
from docx import Document
from pptx import Presentation
from pypdf import PdfReader

from app.api.report import ReportRequest
from app.services.report_service import (
    MAX_LOGO_BYTES,
    _decode_logo,
    build_pdf_report,
    build_powerpoint_report,
    build_word_report,
)


REPORT = {
    "titre": "Comprendre les tendances de fréquentation",
    "synthese": "La fréquentation progresse, mais la qualité inégale de certaines variables limite les comparaisons fines.",
    "messages_cles": [
        "La tendance générale est positive.",
        "Les valeurs manquantes restent concentrées sur deux variables.",
        "La prochaine décision doit prioriser la fiabilisation de la collecte.",
    ],
    "qualite_donnees": {
        "synthese": "Le jeu couvre le périmètre attendu et reste exploitable pour une analyse descriptive.",
        "points_attention": ["Deux variables présentent des valeurs manquantes."],
    },
    "analyses": [
        {
            "titre": "La fréquentation progresse sur la période observée",
            "constat": "Les observations les plus récentes se situent au-dessus du niveau initial.",
            "preuve": "La moyenne passe de 120 à 146 unités.",
            "implication": "La capacité opérationnelle doit être vérifiée avant de prolonger cette tendance.",
        },
        {
            "titre": "La qualité des données conditionne les comparaisons",
            "constat": "Les valeurs absentes ne sont pas réparties uniformément.",
            "preuve": "4,2 % des valeurs sont manquantes.",
            "implication": "Les décisions par segment doivent rester prudentes.",
        },
    ],
    "recommandations": [
        {
            "priorite": "Haute",
            "action": "Corriger les deux variables les plus incomplètes.",
            "justification": "Elles concentrent l'essentiel des valeurs manquantes.",
        },
        {
            "priorite": "Moyenne",
            "action": "Suivre la fréquentation chaque mois.",
            "justification": "La tendance doit être confirmée sur de nouvelles observations.",
        },
    ],
    "limites": ["L'analyse décrit une association et non une causalité."],
    "prochaines_etapes": ["Documenter les règles de collecte.", "Actualiser l'analyse dans trois mois."],
}

PROFILE = {"rows": 1240, "columns": 12}
STATS = {
    "dataset_overview": {
        "pct_valeurs_manquantes_total": 4.2,
        "pct_doublons": 0.8,
    }
}


def _logo_data_url() -> str:
    image = Image.new("RGBA", (320, 100), (102, 117, 73, 255))
    output = io.BytesIO()
    image.save(output, format="PNG")
    return "data:image/png;base64," + base64.b64encode(output.getvalue()).decode("ascii")


def _builder_kwargs() -> dict:
    return {
        "title": "Rapport d'analyse de données",
        "institution": "CITADEL — ancien client à ignorer",
        "filename": "frequentation.csv",
        "analysis_text": "",
        "messages": [],
        "images_b64": [],
        "profile": PROFILE,
        "stats": STATS,
        "models": [],
        "logo_b64": _logo_data_url(),
    }


def _ppt_text(data: bytes) -> str:
    presentation = Presentation(io.BytesIO(data))
    return "\n".join(
        shape.text
        for slide in presentation.slides
        for shape in slide.shapes
        if hasattr(shape, "text") and shape.text
    )


@patch("app.services.report_service.draft_report_with_llm", return_value=REPORT)
def test_exports_use_sali_cover_without_citadel(_draft):
    pdf = build_pdf_report(**_builder_kwargs())
    pdf_text = "\n".join(page.extract_text() or "" for page in PdfReader(io.BytesIO(pdf)).pages)
    assert "Généré par Sali AI" in pdf_text
    assert "CITADEL" not in pdf_text.upper()

    word = build_word_report(**_builder_kwargs())
    document = Document(io.BytesIO(word))
    word_text = "\n".join(paragraph.text for paragraph in document.paragraphs)
    footer_text = "\n".join(
        paragraph.text
        for section in document.sections
        for paragraph in section.first_page_footer.paragraphs
    )
    assert "Généré par Sali AI" in footer_text
    assert "CITADEL" not in (word_text + footer_text).upper()

    powerpoint = build_powerpoint_report(**_builder_kwargs())
    powerpoint_text = _ppt_text(powerpoint)
    assert powerpoint_text.count("Généré par Sali AI") == 1
    assert "CITADEL" not in powerpoint_text.upper()


@patch("app.services.report_service.draft_report_with_llm", return_value=REPORT)
def test_powerpoint_is_a_decision_deck_not_a_text_dump(_draft):
    data = build_powerpoint_report(**_builder_kwargs())
    presentation = Presentation(io.BytesIO(data))
    text = _ppt_text(data)

    assert len(presentation.slides) >= 7
    assert "À retenir" in text
    assert "Le périmètre en un coup d'œil" in text
    assert "Plan d'action" in text
    assert "Limites et prochaines étapes" in text
    assert len(presentation.slides[0].shapes) >= 8  # couverture + logo


def test_logo_validation_and_api_default_are_brand_neutral():
    oversized = base64.b64encode(b"x" * (MAX_LOGO_BYTES + 1)).decode("ascii")
    with pytest.raises(ValueError, match="moins de 3 Mo"):
        _decode_logo(oversized)

    request = ReportRequest(session_id="session-test")
    assert request.institution == ""
    assert request.logo_b64 is None
