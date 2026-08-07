import base64
import html
import io
import json
import logging
import re
from datetime import datetime

from PIL import Image as PILImage
from app.core import config
from app.services.gemini_service import complete_text
from docx import Document
from docx.enum.table import WD_CELL_VERTICAL_ALIGNMENT, WD_TABLE_ALIGNMENT
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml import OxmlElement
from docx.oxml.ns import qn
from docx.shared import Inches, Pt, RGBColor
from pptx import Presentation
from pptx.dml.color import RGBColor as PptxRGB
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches as PptxInches, Pt as PptxPt
from reportlab.lib.colors import HexColor
from reportlab.lib.enums import TA_CENTER, TA_LEFT
from reportlab.lib.pagesizes import A4
from reportlab.lib.styles import ParagraphStyle
from reportlab.lib.units import cm
from reportlab.lib.utils import ImageReader
from reportlab.platypus import (
    HRFlowable,
    Image,
    PageBreak,
    Paragraph,
    SimpleDocTemplate,
    Spacer,
    Table,
    TableStyle,
)

logger = logging.getLogger(__name__)

# Palette éditoriale Sali AI : chaleureuse, sobre et lisible à l'impression.
SALI_OLIVE_HEX = "#667549"
SALI_OLIVE_LIGHT_HEX = "#E8ECDF"
SALI_ORANGE_HEX = "#C67139"
SALI_INK_HEX = "#201E1D"
SALI_MUTED_HEX = "#6F6A64"
SALI_PAPER_HEX = "#F9F4ED"
SALI_WHITE_HEX = "#FFFDFC"
SALI_LINE_HEX = "#D8D0C5"

SALI_OLIVE = HexColor(SALI_OLIVE_HEX)
SALI_OLIVE_LIGHT = HexColor(SALI_OLIVE_LIGHT_HEX)
SALI_ORANGE = HexColor(SALI_ORANGE_HEX)
SALI_INK = HexColor(SALI_INK_HEX)
SALI_MUTED = HexColor(SALI_MUTED_HEX)
SALI_PAPER = HexColor(SALI_PAPER_HEX)
SALI_WHITE = HexColor(SALI_WHITE_HEX)
SALI_LINE = HexColor(SALI_LINE_HEX)

PPTX_OLIVE = PptxRGB(0x66, 0x75, 0x49)
PPTX_OLIVE_LIGHT = PptxRGB(0xE8, 0xEC, 0xDF)
PPTX_ORANGE = PptxRGB(0xC6, 0x71, 0x39)
PPTX_INK = PptxRGB(0x20, 0x1E, 0x1D)
PPTX_MUTED = PptxRGB(0x6F, 0x6A, 0x64)
PPTX_PAPER = PptxRGB(0xF9, 0xF4, 0xED)
PPTX_WHITE = PptxRGB(0xFF, 0xFD, 0xFC)
PPTX_LINE = PptxRGB(0xD8, 0xD0, 0xC5)

MAX_LOGO_BYTES = 3 * 1024 * 1024
MAX_INSIGHTS = 5
MAX_RECOMMENDATIONS = 5


def _clean_text(value) -> str:
    if value is None:
        return ""
    return re.sub(r"[ \t]+", " ", str(value)).strip()


def _pdf_text(value) -> str:
    return html.escape(_clean_text(value)).replace("\n", "<br/>")


def _list_of_text(value, limit: int = 5) -> list[str]:
    if isinstance(value, list):
        values = value
    elif isinstance(value, str):
        values = [
            re.sub(r"^\s*(?:[-•]|\d+[.)])\s*", "", line)
            for line in value.splitlines()
            if line.strip()
        ]
        if not values and value.strip():
            values = [value]
    elif value:
        values = [value]
    else:
        values = []
    return [_clean_text(item) for item in values if _clean_text(item)][:limit]


def _normalise_report(raw: object, fallback_title: str) -> dict:
    """Normalise la réponse LLM vers un contrat utile aux trois formats.

    Le repli sur les anciennes clés permet de rester compatible avec une réponse
    mise en cache ou un modèle qui n'aurait suivi que partiellement le schéma.
    """
    source = raw if isinstance(raw, dict) else {}

    quality_raw = source.get("qualite_donnees") or source.get("description_donnees") or {}
    if isinstance(quality_raw, dict):
        quality = {
            "synthese": _clean_text(quality_raw.get("synthese") or quality_raw.get("constat")),
            "points_attention": _list_of_text(
                quality_raw.get("points_attention") or quality_raw.get("limites"), 5
            ),
        }
    else:
        quality = {"synthese": _clean_text(quality_raw), "points_attention": []}

    analyses_raw = source.get("analyses") or []
    if not analyses_raw and source.get("resultats"):
        analyses_raw = [{"titre": "Résultats principaux", "constat": source.get("resultats")}]
    if isinstance(analyses_raw, dict):
        analyses_raw = [analyses_raw]
    if isinstance(analyses_raw, str):
        analyses_raw = [{"titre": "Résultats principaux", "constat": analyses_raw}]

    analyses = []
    for item in analyses_raw if isinstance(analyses_raw, list) else []:
        if isinstance(item, dict):
            analyses.append({
                "titre": _clean_text(item.get("titre") or "Analyse"),
                "constat": _clean_text(item.get("constat") or item.get("analyse")),
                "preuve": _clean_text(item.get("preuve") or item.get("chiffre_cle")),
                "implication": _clean_text(item.get("implication") or item.get("impact")),
            })
        elif _clean_text(item):
            analyses.append({
                "titre": "Analyse",
                "constat": _clean_text(item),
                "preuve": "",
                "implication": "",
            })
    analyses = analyses[:MAX_INSIGHTS]

    recommendations_raw = source.get("recommandations") or []
    if isinstance(recommendations_raw, dict):
        recommendations_raw = [recommendations_raw]
    if isinstance(recommendations_raw, str):
        recommendations_raw = _list_of_text(recommendations_raw, MAX_RECOMMENDATIONS)

    recommendations = []
    for item in recommendations_raw if isinstance(recommendations_raw, list) else []:
        if isinstance(item, dict):
            priority = _clean_text(item.get("priorite") or "À planifier").capitalize()
            recommendations.append({
                "priorite": priority,
                "action": _clean_text(item.get("action") or item.get("titre")),
                "justification": _clean_text(item.get("justification") or item.get("raison")),
            })
        elif _clean_text(item):
            recommendations.append({
                "priorite": "À planifier",
                "action": _clean_text(item),
                "justification": "",
            })
    recommendations = recommendations[:MAX_RECOMMENDATIONS]

    summary = _clean_text(source.get("synthese") or source.get("resume_executif"))
    key_messages = _list_of_text(source.get("messages_cles"), 4)
    if not key_messages and source.get("conclusions"):
        key_messages = _list_of_text(source.get("conclusions"), 4)

    return {
        "titre": _clean_text(source.get("titre")) or fallback_title,
        "synthese": summary,
        "messages_cles": key_messages,
        "qualite_donnees": quality,
        "analyses": analyses,
        "recommandations": recommendations,
        "limites": _list_of_text(source.get("limites"), 5),
        "prochaines_etapes": _list_of_text(source.get("prochaines_etapes"), 5),
    }


def _decode_logo(logo_b64: str | None) -> bytes | None:
    """Valide le logo et le normalise en PNG pour tous les moteurs d'export."""
    if not logo_b64:
        return None
    payload = logo_b64.split(",", 1)[1] if "," in logo_b64 else logo_b64
    try:
        raw = base64.b64decode(payload, validate=True)
    except Exception as exc:
        raise ValueError("Le logo transmis n'est pas une image valide.") from exc
    if not raw or len(raw) > MAX_LOGO_BYTES:
        raise ValueError("Le logo doit peser moins de 3 Mo.")

    try:
        with PILImage.open(io.BytesIO(raw)) as image:
            image.load()
            image.thumbnail((1800, 1800))
            if image.mode not in ("RGB", "RGBA"):
                image = image.convert("RGBA")
            output = io.BytesIO()
            image.save(output, format="PNG", optimize=True)
            return output.getvalue()
    except Exception as exc:
        raise ValueError("Le format du logo n'est pas pris en charge.") from exc


def _logo_ratio(logo_bytes: bytes) -> float:
    with PILImage.open(io.BytesIO(logo_bytes)) as image:
        width, height = image.size
    return width / max(height, 1)


def _format_number(value) -> str:
    if value is None or value == "":
        return ""
    try:
        number = float(value)
    except (TypeError, ValueError):
        return _clean_text(value)
    if number.is_integer():
        return f"{int(number):,}".replace(",", " ")
    return f"{number:,.1f}".replace(",", " ").replace(".", ",")


def _report_metrics(profile: dict | None, stats: dict | None) -> list[dict[str, str]]:
    profile = profile or {}
    overview = (stats or {}).get("dataset_overview", {}) or {}
    rows = profile.get("rows", overview.get("n_lignes"))
    columns = profile.get("columns", overview.get("n_variables"))
    missing = overview.get("pct_valeurs_manquantes_total")
    duplicates = overview.get("pct_doublons")

    metrics = []
    if rows is not None:
        metrics.append({"label": "Observations", "value": _format_number(rows)})
    if columns is not None:
        metrics.append({"label": "Variables", "value": _format_number(columns)})
    if missing is not None:
        metrics.append({"label": "Valeurs manquantes", "value": f"{_format_number(missing)} %"})
    if duplicates is not None:
        metrics.append({"label": "Doublons", "value": f"{_format_number(duplicates)} %"})
    return metrics[:4]


def _bloc_faits(profile: dict, stats: dict, models: list) -> str:
    """Faits calculés transmis au LLM : seule source autorisée pour les nombres."""
    profile = profile or {}
    stats = stats or {}
    overview = stats.get("dataset_overview", {}) or {}
    variables = stats.get("variables", {}) or {}

    lines = [
        "FAITS VÉRIFIÉS — reprendre les valeurs telles quelles",
        f"- Lignes : {profile.get('rows', overview.get('n_lignes', 'inconnu'))}",
        f"- Colonnes : {profile.get('columns', overview.get('n_variables', 'inconnu'))}",
    ]
    if profile.get("column_names"):
        lines.append(f"- Colonnes : {', '.join(map(str, profile['column_names']))}")
    if overview:
        lines.extend([
            f"- Doublons : {overview.get('n_doublons', 0)} ({overview.get('pct_doublons', 0)} %)",
            "- Valeurs manquantes : "
            f"{overview.get('n_valeurs_manquantes_total', 0)} "
            f"({overview.get('pct_valeurs_manquantes_total', 0)} %)",
            "- Variables numériques : "
            f"{overview.get('n_variables_numeriques', 0)} ; catégorielles : "
            f"{overview.get('n_variables_categorielles', 0)}",
        ])

    compact_variables = {}
    allowed = {
        "type", "moyenne", "mediane", "ecart_type", "min", "max", "q1", "q3",
        "skewness", "n_manquantes", "n_valeurs_distinctes", "valeur_dominante",
        "frequence_dominante",
    }
    for column, values in variables.items():
        if isinstance(values, dict):
            compact_variables[column] = {key: values[key] for key in allowed if key in values}
    if compact_variables:
        lines.append("\nSTATISTIQUES PAR VARIABLE :")
        lines.append(json.dumps(compact_variables, ensure_ascii=False, indent=2))

    for label, key in (
        ("CORRÉLATIONS", "correlations"),
        ("VALEURS MANQUANTES PAR COLONNE", "missing"),
    ):
        value = stats.get(key)
        if value:
            lines.append(f"\n{label} :")
            lines.append(json.dumps(value, ensure_ascii=False, indent=2))

    if models:
        lines.append("\nMODÈLES ENTRAÎNÉS :")
        for model in models:
            features = ", ".join(map(str, model.get("features") or [])) or "non précisées"
            lines.append(
                f"- {model.get('name')} — type : {model.get('type')} ; variables : {features}"
            )
            if model.get("metrics"):
                lines.append(f"  métriques : {json.dumps(model['metrics'], ensure_ascii=False)}")
    return "\n".join(lines)


def draft_report_with_llm(
    filename: str,
    analysis_text: str,
    chat_history: list,
    title: str,
    institution: str = "",
    model: str | None = None,
    key_points: str = "",
    profile: dict | None = None,
    stats: dict | None = None,
    models: list | None = None,
) -> dict:
    """Rédige un contenu décisionnel commun, ensuite adapté au document ou au deck."""
    del institution  # Conservé dans l'API pour compatibilité, jamais injecté dans l'export.
    model = model or config.get_default_model()
    history_summary = "\n".join(
        f"[{message.get('role', '').upper()}] {_clean_text(message.get('text'))[:700]}"
        for message in chat_history
        if isinstance(message, dict)
    )
    commissioner = ""
    if key_points.strip():
        commissioner = f"""
DEMANDE PRIORITAIRE DU COMMANDITAIRE :
{key_points.strip()}

Cette demande détermine l'angle du rapport. Si les données n'y répondent pas,
signale explicitement la limite au lieu de la contourner.
"""

    prompt = f"""
Tu es analyste senior et rédacteur de livrables de décision.

SOURCE ANALYSÉE : {filename}
TITRE DE REPLI : {title}
{commissioner}
{_bloc_faits(profile or {}, stats or {}, models or [])}

ANALYSE INITIALE :
{analysis_text}

ÉCHANGES UTILES DE LA SESSION :
{history_summary}

Produis un rapport en français, précis et directement exploitable par un décideur.
Sépare toujours un fait observé de son interprétation. Chaque nombre doit venir
du bloc FAITS VÉRIFIÉS. N'invente ni causalité, ni période, ni objectif métier.
Privilégie les conclusions spécifiques à la source aux formulations génériques.

Retourne EXCLUSIVEMENT ce JSON valide, sans markdown :
{{
  "titre": "titre spécifique de 6 à 12 mots",
  "synthese": "80 à 120 mots : réponse directe, portée et enjeu principal",
  "messages_cles": ["3 ou 4 messages autonomes, chacun en 20 mots maximum"],
  "qualite_donnees": {{
    "synthese": "qualité, couverture et aptitude à l'analyse en 80 mots maximum",
    "points_attention": ["limites de qualité réellement observées"]
  }},
  "analyses": [
    {{
      "titre": "conclusion analytique formulée comme un message",
      "constat": "ce qui est observé",
      "preuve": "le ou les chiffres vérifiés qui l'étayent",
      "implication": "ce que cela change pour la décision, sans inventer le contexte"
    }}
  ],
  "recommandations": [
    {{
      "priorite": "Haute, Moyenne ou À surveiller",
      "action": "action concrète et vérifiable",
      "justification": "lien explicite avec un constat"
    }}
  ],
  "limites": ["ce que les données ou la méthode ne permettent pas d'affirmer"],
  "prochaines_etapes": ["étape réaliste permettant d'approfondir ou de décider"]
}}

Contraintes : 3 à 5 analyses, 3 à 5 recommandations, aucune section de remplissage.
"""
    text = complete_text(prompt, model).strip()
    if text.startswith("```"):
        lines = text.splitlines()
        text = "\n".join(lines[1:-1] if lines[-1].strip().startswith("```") else lines[1:])
    try:
        raw = json.loads(text)
    except json.JSONDecodeError:
        logger.warning("Rapport : réponse non-JSON du modèle %s, repli sur le texte brut.", model)
        raw = {"synthese": text}
    return _normalise_report(raw, title)


def _fit_dimensions(width: float, height: float, max_width: float, max_height: float) -> tuple[float, float]:
    ratio = min(max_width / max(width, 1), max_height / max(height, 1))
    return width * ratio, height * ratio


def _image_dimensions(raw: bytes, max_width: float, max_height: float) -> tuple[float, float]:
    with PILImage.open(io.BytesIO(raw)) as image:
        width, height = image.size
    return _fit_dimensions(width, height, max_width, max_height)


def _draft(
    *, title, institution, filename, analysis_text, messages, model, key_points, profile, stats, models
) -> dict:
    return draft_report_with_llm(
        filename=filename,
        analysis_text=analysis_text,
        chat_history=messages,
        title=title,
        institution=institution,
        model=model,
        key_points=key_points,
        profile=profile,
        stats=stats,
        models=models,
    )


def build_pdf_report(
    title: str,
    institution: str,
    filename: str,
    analysis_text: str,
    messages: list,
    images_b64: list,
    model: str | None = None,
    key_points: str = "",
    profile: dict | None = None,
    stats: dict | None = None,
    models: list | None = None,
    logo_b64: str | None = None,
) -> bytes:
    logo = _decode_logo(logo_b64)
    report = _draft(
        title=title, institution=institution, filename=filename,
        analysis_text=analysis_text, messages=messages, model=model,
        key_points=key_points, profile=profile, stats=stats, models=models,
    )
    metrics = _report_metrics(profile, stats)

    buffer = io.BytesIO()
    document = SimpleDocTemplate(
        buffer,
        pagesize=A4,
        rightMargin=2.2 * cm,
        leftMargin=2.2 * cm,
        topMargin=2.0 * cm,
        bottomMargin=1.8 * cm,
        title=report["titre"],
        author="Sali AI",
    )
    page_width, page_height = A4

    title_style = ParagraphStyle(
        "ReportTitle", fontName="Helvetica-Bold", fontSize=28, leading=33,
        textColor=SALI_INK, alignment=TA_LEFT,
    )
    section_style = ParagraphStyle(
        "Section", fontName="Helvetica-Bold", fontSize=17, leading=21,
        textColor=SALI_OLIVE, spaceBefore=18, spaceAfter=9,
    )
    insight_style = ParagraphStyle(
        "Insight", fontName="Helvetica-Bold", fontSize=12.5, leading=16,
        textColor=SALI_INK, spaceBefore=12, spaceAfter=5,
    )
    body_style = ParagraphStyle(
        "Body", fontName="Helvetica", fontSize=10.2, leading=15.5,
        textColor=SALI_INK, spaceAfter=7,
    )
    muted_style = ParagraphStyle(
        "Muted", fontName="Helvetica", fontSize=8.5, leading=11,
        textColor=SALI_MUTED,
    )
    label_style = ParagraphStyle(
        "Label", fontName="Helvetica-Bold", fontSize=8, leading=10,
        textColor=SALI_ORANGE, spaceBefore=5, spaceAfter=2,
    )
    bullet_style = ParagraphStyle(
        "Bullet", parent=body_style, leftIndent=14, firstLineIndent=-10,
        bulletIndent=0, spaceAfter=5,
    )

    def draw_cover(canvas, _doc):
        canvas.saveState()
        canvas.setFillColor(SALI_PAPER)
        canvas.rect(0, 0, page_width, page_height, stroke=0, fill=1)
        canvas.setFillColor(SALI_OLIVE)
        canvas.rect(0, page_height - 0.35 * cm, page_width, 0.35 * cm, stroke=0, fill=1)
        canvas.setFillColor(SALI_ORANGE)
        canvas.circle(2.25 * cm, page_height - 2.1 * cm, 0.13 * cm, stroke=0, fill=1)

        if logo:
            ratio = _logo_ratio(logo)
            max_width, max_height = 4.0 * cm, 1.8 * cm
            width, height = _fit_dimensions(ratio, 1, max_width, max_height)
            canvas.drawImage(
                ImageReader(io.BytesIO(logo)), page_width - 2.2 * cm - width,
                page_height - 2.8 * cm, width=width, height=height,
                preserveAspectRatio=True, mask="auto",
            )

        canvas.setFillColor(SALI_OLIVE)
        canvas.setFont("Helvetica-Bold", 9)
        canvas.drawString(2.2 * cm, page_height - 2.3 * cm, "RAPPORT D'ANALYSE")
        cover_title = Paragraph(_pdf_text(report["titre"]), title_style)
        _, title_height = cover_title.wrap(page_width - 4.4 * cm, 7 * cm)
        cover_title.drawOn(canvas, 2.2 * cm, page_height - 5.0 * cm - title_height)

        canvas.setStrokeColor(SALI_ORANGE)
        canvas.setLineWidth(2.2)
        canvas.line(2.2 * cm, page_height - 6.0 * cm - title_height, 5.2 * cm, page_height - 6.0 * cm - title_height)
        canvas.setFillColor(SALI_MUTED)
        canvas.setFont("Helvetica", 10)
        source = _clean_text(filename) or "Source de la session"
        canvas.drawString(2.2 * cm, 5.4 * cm, f"Source analysée · {source[:85]}")
        canvas.drawString(2.2 * cm, 4.8 * cm, datetime.now().strftime("Édition du %d/%m/%Y"))

        canvas.setStrokeColor(SALI_LINE)
        canvas.setLineWidth(0.7)
        canvas.line(2.2 * cm, 2.0 * cm, page_width - 2.2 * cm, 2.0 * cm)
        canvas.setFillColor(SALI_MUTED)
        canvas.setFont("Helvetica", 8.5)
        canvas.drawString(2.2 * cm, 1.45 * cm, "Généré par Sali AI")
        canvas.restoreState()

    def draw_page_number(canvas, doc):
        canvas.saveState()
        canvas.setFillColor(SALI_MUTED)
        canvas.setFont("Helvetica", 8)
        canvas.drawRightString(page_width - 2.2 * cm, 0.9 * cm, str(doc.page))
        canvas.restoreState()

    def add_section(elements: list, heading: str):
        elements.append(Paragraph(_pdf_text(heading), section_style))
        elements.append(HRFlowable(width="100%", thickness=0.7, color=SALI_LINE, spaceAfter=8))

    elements: list = [PageBreak()]
    add_section(elements, "Synthèse décisionnelle")
    if report["synthese"]:
        summary = Table(
            [[Paragraph(_pdf_text(report["synthese"]), body_style)]],
            colWidths=[16.6 * cm],
        )
        summary.setStyle(TableStyle([
            ("BACKGROUND", (0, 0), (-1, -1), SALI_PAPER),
            ("BOX", (0, 0), (-1, -1), 0.8, SALI_LINE),
            ("LEFTPADDING", (0, 0), (-1, -1), 14),
            ("RIGHTPADDING", (0, 0), (-1, -1), 14),
            ("TOPPADDING", (0, 0), (-1, -1), 12),
            ("BOTTOMPADDING", (0, 0), (-1, -1), 10),
        ]))
        elements.extend([summary, Spacer(1, 10)])
    for index, message in enumerate(report["messages_cles"], 1):
        elements.append(Paragraph(f"<b>{index:02d}</b> &nbsp; {_pdf_text(message)}", bullet_style))

    if metrics:
        add_section(elements, "Les données en un coup d'œil")
        cells = []
        for metric in metrics:
            cells.append(Paragraph(
                f"<font size='17' color='{SALI_INK_HEX}'><b>{_pdf_text(metric['value'])}</b></font>"
                f"<br/><font size='8' color='{SALI_MUTED_HEX}'>{_pdf_text(metric['label']).upper()}</font>",
                ParagraphStyle("Metric", alignment=TA_CENTER, leading=18),
            ))
        metric_table = Table([cells], colWidths=[16.6 * cm / len(cells)] * len(cells))
        metric_table.setStyle(TableStyle([
            ("BACKGROUND", (0, 0), (-1, -1), SALI_WHITE),
            ("GRID", (0, 0), (-1, -1), 0.7, SALI_LINE),
            ("TOPPADDING", (0, 0), (-1, -1), 13),
            ("BOTTOMPADDING", (0, 0), (-1, -1), 13),
            ("VALIGN", (0, 0), (-1, -1), "MIDDLE"),
        ]))
        elements.append(metric_table)

    quality = report["qualite_donnees"]
    if quality["synthese"] or quality["points_attention"]:
        add_section(elements, "Périmètre et qualité des données")
        if quality["synthese"]:
            elements.append(Paragraph(_pdf_text(quality["synthese"]), body_style))
        for point in quality["points_attention"]:
            elements.append(Paragraph(f"• {_pdf_text(point)}", bullet_style))

    if report["analyses"]:
        add_section(elements, "Analyses et implications")
        for index, insight in enumerate(report["analyses"], 1):
            elements.append(Paragraph(
                f"{index}. {_pdf_text(insight['titre'])}", insight_style
            ))
            if insight["constat"]:
                elements.append(Paragraph("CONSTAT", label_style))
                elements.append(Paragraph(_pdf_text(insight["constat"]), body_style))
            if insight["preuve"]:
                evidence = Table(
                    [[Paragraph(f"<b>Preuve</b><br/>{_pdf_text(insight['preuve'])}", body_style)]],
                    colWidths=[16.1 * cm],
                )
                evidence.setStyle(TableStyle([
                    ("BACKGROUND", (0, 0), (-1, -1), SALI_OLIVE_LIGHT),
                    ("LINEBEFORE", (0, 0), (0, -1), 3, SALI_OLIVE),
                    ("LEFTPADDING", (0, 0), (-1, -1), 10),
                    ("RIGHTPADDING", (0, 0), (-1, -1), 10),
                    ("TOPPADDING", (0, 0), (-1, -1), 8),
                    ("BOTTOMPADDING", (0, 0), (-1, -1), 7),
                ]))
                elements.extend([evidence, Spacer(1, 4)])
            if insight["implication"]:
                elements.append(Paragraph("IMPLICATION", label_style))
                elements.append(Paragraph(_pdf_text(insight["implication"]), body_style))

    valid_images = []
    for value in images_b64 or []:
        try:
            valid_images.append(base64.b64decode(value))
        except Exception:
            logger.warning("Une visualisation illisible a été ignorée dans le PDF.")
    if valid_images:
        add_section(elements, "Visualisations")
        for index, raw in enumerate(valid_images, 1):
            try:
                width, height = _image_dimensions(raw, 15.8 * cm, 9.5 * cm)
                figure = Image(io.BytesIO(raw), width=width, height=height)
                figure.hAlign = "CENTER"
                elements.extend([
                    figure,
                    Paragraph(f"Figure {index}", ParagraphStyle("Caption", parent=muted_style, alignment=TA_CENTER)),
                    Spacer(1, 10),
                ])
            except Exception as exc:
                logger.warning("Visualisation %s non insérée dans le PDF : %s", index, exc)

    if report["recommandations"]:
        add_section(elements, "Plan d'action recommandé")
        rows = [[
            Paragraph("PRIORITÉ", label_style),
            Paragraph("ACTION", label_style),
            Paragraph("POURQUOI", label_style),
        ]]
        for recommendation in report["recommandations"]:
            rows.append([
                Paragraph(_pdf_text(recommendation["priorite"]), body_style),
                Paragraph(_pdf_text(recommendation["action"]), body_style),
                Paragraph(_pdf_text(recommendation["justification"]), body_style),
            ])
        action_table = Table(rows, colWidths=[2.4 * cm, 6.1 * cm, 8.1 * cm], repeatRows=1)
        action_table.setStyle(TableStyle([
            ("BACKGROUND", (0, 0), (-1, 0), SALI_PAPER),
            ("GRID", (0, 0), (-1, -1), 0.6, SALI_LINE),
            ("VALIGN", (0, 0), (-1, -1), "TOP"),
            ("LEFTPADDING", (0, 0), (-1, -1), 7),
            ("RIGHTPADDING", (0, 0), (-1, -1), 7),
            ("TOPPADDING", (0, 0), (-1, -1), 7),
            ("BOTTOMPADDING", (0, 0), (-1, -1), 6),
        ]))
        elements.append(action_table)

    if report["limites"] or report["prochaines_etapes"]:
        add_section(elements, "Limites et prochaines étapes")
        if report["limites"]:
            elements.append(Paragraph("LIMITES", label_style))
            for item in report["limites"]:
                elements.append(Paragraph(f"• {_pdf_text(item)}", bullet_style))
        if report["prochaines_etapes"]:
            elements.append(Paragraph("PROCHAINES ÉTAPES", label_style))
            for item in report["prochaines_etapes"]:
                elements.append(Paragraph(f"• {_pdf_text(item)}", bullet_style))

    document.build(elements, onFirstPage=draw_cover, onLaterPages=draw_page_number)
    buffer.seek(0)
    return buffer.read()


def _shade_word_cell(cell, fill: str) -> None:
    properties = cell._tc.get_or_add_tcPr()
    shading = OxmlElement("w:shd")
    shading.set(qn("w:fill"), fill)
    properties.append(shading)


def _set_word_cell_text(cell, text: str, *, bold: bool = False, size: float = 9.5, color: str = "201E1D") -> None:
    cell.text = ""
    paragraph = cell.paragraphs[0]
    run = paragraph.add_run(text)
    run.bold = bold
    run.font.name = "Aptos"
    run.font.size = Pt(size)
    run.font.color.rgb = RGBColor.from_string(color)
    cell.vertical_alignment = WD_CELL_VERTICAL_ALIGNMENT.CENTER


def build_word_report(
    title: str,
    institution: str,
    filename: str,
    analysis_text: str,
    messages: list,
    images_b64: list,
    model: str | None = None,
    key_points: str = "",
    profile: dict | None = None,
    stats: dict | None = None,
    models: list | None = None,
    logo_b64: str | None = None,
) -> bytes:
    logo = _decode_logo(logo_b64)
    report = _draft(
        title=title, institution=institution, filename=filename,
        analysis_text=analysis_text, messages=messages, model=model,
        key_points=key_points, profile=profile, stats=stats, models=models,
    )
    metrics = _report_metrics(profile, stats)
    document = Document()
    section = document.sections[0]
    section.top_margin = Inches(0.85)
    section.bottom_margin = Inches(0.8)
    section.left_margin = Inches(1.0)
    section.right_margin = Inches(1.0)
    section.different_first_page_header_footer = True

    normal = document.styles["Normal"]
    normal.font.name = "Aptos"
    normal.font.size = Pt(10.5)
    normal.font.color.rgb = RGBColor.from_string("201E1D")
    normal.paragraph_format.space_after = Pt(6)

    for name, size, color in (
        ("Title", 27, "201E1D"),
        ("Heading 1", 17, "667549"),
        ("Heading 2", 12.5, "201E1D"),
    ):
        style = document.styles[name]
        style.font.name = "Aptos Display"
        style.font.size = Pt(size)
        style.font.color.rgb = RGBColor.from_string(color)
        style.font.bold = True

    footer = section.first_page_footer.paragraphs[0]
    footer.text = "Généré par Sali AI"
    footer.alignment = WD_ALIGN_PARAGRAPH.LEFT
    footer.runs[0].font.name = "Aptos"
    footer.runs[0].font.size = Pt(8.5)
    footer.runs[0].font.color.rgb = RGBColor.from_string("6F6A64")

    if logo:
        document.add_picture(io.BytesIO(logo), width=Inches(1.55))
        document.paragraphs[-1].alignment = WD_ALIGN_PARAGRAPH.RIGHT
    eyebrow = document.add_paragraph()
    run = eyebrow.add_run("RAPPORT D'ANALYSE")
    run.bold = True
    run.font.name = "Aptos"
    run.font.size = Pt(9)
    run.font.color.rgb = RGBColor.from_string("667549")
    eyebrow.paragraph_format.space_before = Pt(60 if logo else 90)
    eyebrow.paragraph_format.space_after = Pt(12)
    document.add_paragraph(report["titre"], style="Title")
    accent = document.add_paragraph("━━━")
    accent.runs[0].font.color.rgb = RGBColor.from_string("C67139")
    accent.runs[0].font.size = Pt(13)
    accent.paragraph_format.space_after = Pt(22)
    source = document.add_paragraph(f"Source analysée · {_clean_text(filename) or 'Source de la session'}")
    source.runs[0].font.color.rgb = RGBColor.from_string("6F6A64")
    source.add_run(f"\nÉdition du {datetime.now().strftime('%d/%m/%Y')}")
    source.runs[-1].font.color.rgb = RGBColor.from_string("6F6A64")
    document.add_page_break()

    def heading(text: str):
        paragraph = document.add_paragraph(text, style="Heading 1")
        paragraph.paragraph_format.space_before = Pt(12)
        paragraph.paragraph_format.space_after = Pt(7)
        return paragraph

    def bullets(items: list[str]):
        for item in items:
            paragraph = document.add_paragraph(style="List Bullet")
            paragraph.add_run(item)

    heading("Synthèse décisionnelle")
    if report["synthese"]:
        table = document.add_table(rows=1, cols=1)
        table.alignment = WD_TABLE_ALIGNMENT.CENTER
        _shade_word_cell(table.cell(0, 0), "F9F4ED")
        _set_word_cell_text(table.cell(0, 0), report["synthese"], size=10.5)
        document.add_paragraph()
    for index, message in enumerate(report["messages_cles"], 1):
        paragraph = document.add_paragraph()
        number = paragraph.add_run(f"{index:02d}  ")
        number.bold = True
        number.font.color.rgb = RGBColor.from_string("C67139")
        paragraph.add_run(message)

    if metrics:
        heading("Les données en un coup d'œil")
        table = document.add_table(rows=1, cols=len(metrics))
        table.alignment = WD_TABLE_ALIGNMENT.CENTER
        for index, metric in enumerate(metrics):
            cell = table.cell(0, index)
            _shade_word_cell(cell, "F9F4ED")
            cell.text = ""
            paragraph = cell.paragraphs[0]
            paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
            value = paragraph.add_run(metric["value"])
            value.bold = True
            value.font.name = "Aptos Display"
            value.font.size = Pt(17)
            value.font.color.rgb = RGBColor.from_string("201E1D")
            label = paragraph.add_run(f"\n{metric['label'].upper()}")
            label.font.name = "Aptos"
            label.font.size = Pt(7.5)
            label.font.color.rgb = RGBColor.from_string("6F6A64")

    quality = report["qualite_donnees"]
    if quality["synthese"] or quality["points_attention"]:
        heading("Périmètre et qualité des données")
        if quality["synthese"]:
            document.add_paragraph(quality["synthese"])
        bullets(quality["points_attention"])

    if report["analyses"]:
        heading("Analyses et implications")
        for index, insight in enumerate(report["analyses"], 1):
            document.add_paragraph(f"{index}. {insight['titre']}", style="Heading 2")
            if insight["constat"]:
                document.add_paragraph(insight["constat"])
            if insight["preuve"]:
                table = document.add_table(rows=1, cols=1)
                _shade_word_cell(table.cell(0, 0), "E8ECDF")
                _set_word_cell_text(table.cell(0, 0), f"PREUVE\n{insight['preuve']}", size=9.5)
            if insight["implication"]:
                paragraph = document.add_paragraph()
                label = paragraph.add_run("IMPLICATION  ")
                label.bold = True
                label.font.color.rgb = RGBColor.from_string("C67139")
                paragraph.add_run(insight["implication"])

    valid_images = []
    for value in images_b64 or []:
        try:
            valid_images.append(base64.b64decode(value))
        except Exception:
            logger.warning("Une visualisation illisible a été ignorée dans Word.")
    if valid_images:
        heading("Visualisations")
        for index, raw in enumerate(valid_images, 1):
            try:
                document.add_picture(io.BytesIO(raw), width=Inches(6.1))
                document.paragraphs[-1].alignment = WD_ALIGN_PARAGRAPH.CENTER
                caption = document.add_paragraph(f"Figure {index}")
                caption.alignment = WD_ALIGN_PARAGRAPH.CENTER
                caption.runs[0].font.size = Pt(8.5)
                caption.runs[0].font.color.rgb = RGBColor.from_string("6F6A64")
            except Exception as exc:
                logger.warning("Visualisation %s non insérée dans Word : %s", index, exc)

    if report["recommandations"]:
        heading("Plan d'action recommandé")
        table = document.add_table(rows=1, cols=3)
        table.alignment = WD_TABLE_ALIGNMENT.CENTER
        table.style = "Table Grid"
        for cell, value in zip(table.rows[0].cells, ("PRIORITÉ", "ACTION", "POURQUOI")):
            _shade_word_cell(cell, "F9F4ED")
            _set_word_cell_text(cell, value, bold=True, size=8, color="C67139")
        for recommendation in report["recommandations"]:
            cells = table.add_row().cells
            _set_word_cell_text(cells[0], recommendation["priorite"], bold=True, size=9)
            _set_word_cell_text(cells[1], recommendation["action"], size=9)
            _set_word_cell_text(cells[2], recommendation["justification"], size=9)

    if report["limites"] or report["prochaines_etapes"]:
        heading("Limites et prochaines étapes")
        if report["limites"]:
            document.add_paragraph("Limites", style="Heading 2")
            bullets(report["limites"])
        if report["prochaines_etapes"]:
            document.add_paragraph("Prochaines étapes", style="Heading 2")
            bullets(report["prochaines_etapes"])

    buffer = io.BytesIO()
    document.save(buffer)
    buffer.seek(0)
    return buffer.read()


def _clip(value: str, maximum: int) -> str:
    text = _clean_text(value)
    if len(text) <= maximum:
        return text
    shortened = text[: maximum - 1].rsplit(" ", 1)[0]
    return f"{shortened}…"


def build_powerpoint_report(
    title: str,
    institution: str,
    filename: str,
    analysis_text: str,
    messages: list,
    images_b64: list,
    model: str | None = None,
    key_points: str = "",
    profile: dict | None = None,
    stats: dict | None = None,
    models: list | None = None,
    logo_b64: str | None = None,
) -> bytes:
    logo = _decode_logo(logo_b64)
    report = _draft(
        title=title, institution=institution, filename=filename,
        analysis_text=analysis_text, messages=messages, model=model,
        key_points=key_points, profile=profile, stats=stats, models=models,
    )
    metrics = _report_metrics(profile, stats)

    presentation = Presentation()
    presentation.slide_width = PptxInches(13.333)
    presentation.slide_height = PptxInches(7.5)
    width = presentation.slide_width
    height = presentation.slide_height
    blank = presentation.slide_layouts[6]
    slide_number = 0

    def rectangle(slide, x, y, w, h, fill, line=None, radius=True):
        shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
        shape = slide.shapes.add_shape(shape_type, x, y, w, h)
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
        if line:
            shape.line.color.rgb = line
            shape.line.width = PptxPt(0.8)
        else:
            shape.line.fill.background()
        shape.shadow.inherit = False
        return shape

    def text_box(
        slide, x, y, w, h, text, size, color, *, bold=False,
        align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP, margin=0,
    ):
        shape = slide.shapes.add_textbox(x, y, w, h)
        frame = shape.text_frame
        frame.clear()
        frame.word_wrap = True
        frame.margin_left = frame.margin_right = frame.margin_top = frame.margin_bottom = margin
        frame.vertical_anchor = valign
        paragraph = frame.paragraphs[0]
        paragraph.text = _clean_text(text)
        paragraph.alignment = align
        paragraph.space_after = PptxPt(0)
        for run in paragraph.runs:
            run.font.name = "Aptos"
            run.font.size = PptxPt(size)
            run.font.color.rgb = color
            run.font.bold = bold
        return shape

    def base_slide(kicker: str, heading: str):
        nonlocal slide_number
        slide_number += 1
        slide = presentation.slides.add_slide(blank)
        rectangle(slide, 0, 0, width, height, PPTX_WHITE, radius=False)
        rectangle(slide, 0, 0, PptxInches(0.16), height, PPTX_OLIVE, radius=False)
        text_box(slide, PptxInches(0.8), PptxInches(0.45), PptxInches(8.5), PptxInches(0.28), kicker.upper(), 9, PPTX_ORANGE, bold=True)
        text_box(slide, PptxInches(0.8), PptxInches(0.82), PptxInches(11.6), PptxInches(0.65), heading, 27, PPTX_INK, bold=True)
        rectangle(slide, PptxInches(0.8), PptxInches(1.52), PptxInches(11.75), Emu(11000), PPTX_LINE, radius=False)
        text_box(slide, PptxInches(12.05), PptxInches(7.02), PptxInches(0.55), PptxInches(0.2), f"{slide_number:02d}", 8, PPTX_MUTED, align=PP_ALIGN.RIGHT)
        return slide

    # Couverture : l'attribution Sali AI n'apparaît qu'ici.
    slide_number += 1
    cover = presentation.slides.add_slide(blank)
    rectangle(cover, 0, 0, width, height, PPTX_PAPER, radius=False)
    rectangle(cover, 0, 0, PptxInches(0.22), height, PPTX_OLIVE, radius=False)
    rectangle(cover, PptxInches(0.8), PptxInches(0.82), PptxInches(0.12), PptxInches(0.12), PPTX_ORANGE)
    text_box(cover, PptxInches(1.04), PptxInches(0.79), PptxInches(4.0), PptxInches(0.25), "RAPPORT D'ANALYSE", 10, PPTX_OLIVE, bold=True)
    text_box(cover, PptxInches(0.85), PptxInches(2.0), PptxInches(10.7), PptxInches(1.55), report["titre"], 37, PPTX_INK, bold=True, valign=MSO_ANCHOR.MIDDLE)
    rectangle(cover, PptxInches(0.86), PptxInches(3.78), PptxInches(1.55), Emu(22000), PPTX_ORANGE, radius=False)
    text_box(
        cover, PptxInches(0.85), PptxInches(4.18), PptxInches(9.5), PptxInches(0.32),
        f"Source analysée · {_clean_text(filename) or 'Source de la session'}", 12, PPTX_MUTED,
    )
    text_box(cover, PptxInches(0.85), PptxInches(4.62), PptxInches(4.0), PptxInches(0.25), datetime.now().strftime("Édition du %d/%m/%Y"), 10, PPTX_MUTED)
    if logo:
        ratio = _logo_ratio(logo)
        max_w, max_h = PptxInches(2.0), PptxInches(0.95)
        logo_w, logo_h = _fit_dimensions(ratio, 1, max_w, max_h)
        cover.shapes.add_picture(
            io.BytesIO(logo), width - PptxInches(0.85) - int(logo_w), PptxInches(0.62),
            width=int(logo_w), height=int(logo_h),
        )
    rectangle(cover, PptxInches(0.85), PptxInches(6.82), PptxInches(11.65), Emu(8500), PPTX_LINE, radius=False)
    text_box(cover, PptxInches(0.85), PptxInches(6.98), PptxInches(3.0), PptxInches(0.2), "Généré par Sali AI", 8.5, PPTX_MUTED)

    if report["messages_cles"] or report["synthese"]:
        slide = base_slide("Synthèse", "À retenir")
        messages_to_show = report["messages_cles"][:4]
        if not messages_to_show and report["synthese"]:
            messages_to_show = [report["synthese"]]
        count = len(messages_to_show)
        card_h = min(1.08, 4.7 / max(count, 1))
        y = 1.86
        for index, message in enumerate(messages_to_show, 1):
            rectangle(slide, PptxInches(0.82), PptxInches(y), PptxInches(11.65), PptxInches(card_h), PPTX_PAPER, PPTX_LINE)
            text_box(slide, PptxInches(1.08), PptxInches(y + 0.18), PptxInches(0.5), PptxInches(0.45), f"{index:02d}", 16, PPTX_ORANGE, bold=True, valign=MSO_ANCHOR.MIDDLE)
            text_box(slide, PptxInches(1.72), PptxInches(y + 0.14), PptxInches(10.2), PptxInches(card_h - 0.24), _clip(message, 210), 17, PPTX_INK, bold=True, valign=MSO_ANCHOR.MIDDLE)
            y += card_h + 0.18

    if metrics:
        slide = base_slide("Données", "Le périmètre en un coup d'œil")
        count = len(metrics)
        gap = 0.22
        card_width = (11.65 - gap * (count - 1)) / count
        for index, metric in enumerate(metrics):
            x = 0.82 + index * (card_width + gap)
            rectangle(slide, PptxInches(x), PptxInches(2.05), PptxInches(card_width), PptxInches(2.55), PPTX_PAPER, PPTX_LINE)
            text_box(slide, PptxInches(x + 0.18), PptxInches(2.48), PptxInches(card_width - 0.36), PptxInches(0.9), metric["value"], 28, PPTX_INK, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
            text_box(slide, PptxInches(x + 0.18), PptxInches(3.58), PptxInches(card_width - 0.36), PptxInches(0.4), metric["label"].upper(), 9, PPTX_OLIVE, bold=True, align=PP_ALIGN.CENTER)
        quality = report["qualite_donnees"]
        quality_text = quality["synthese"] or "Les indicateurs ci-dessus décrivent le périmètre disponible pour l'analyse."
        text_box(slide, PptxInches(1.25), PptxInches(5.14), PptxInches(10.8), PptxInches(0.9), _clip(quality_text, 360), 14, PPTX_MUTED, align=PP_ALIGN.CENTER)

    for index, insight in enumerate(report["analyses"], 1):
        slide = base_slide(f"Analyse {index:02d}", _clip(insight["titre"], 92))
        text_box(slide, PptxInches(0.85), PptxInches(1.86), PptxInches(11.55), PptxInches(1.25), _clip(insight["constat"], 430), 18, PPTX_INK, bold=True, valign=MSO_ANCHOR.MIDDLE)
        rectangle(slide, PptxInches(0.85), PptxInches(3.42), PptxInches(5.55), PptxInches(2.25), PPTX_OLIVE_LIGHT, PPTX_OLIVE)
        text_box(slide, PptxInches(1.12), PptxInches(3.72), PptxInches(4.95), PptxInches(0.25), "PREUVE", 9, PPTX_OLIVE, bold=True)
        text_box(slide, PptxInches(1.12), PptxInches(4.12), PptxInches(4.95), PptxInches(1.16), _clip(insight["preuve"] or "Aucun chiffre supplémentaire disponible.", 320), 15, PPTX_INK, bold=True, valign=MSO_ANCHOR.MIDDLE)
        rectangle(slide, PptxInches(6.68), PptxInches(3.42), PptxInches(5.72), PptxInches(2.25), PPTX_PAPER, PPTX_LINE)
        text_box(slide, PptxInches(6.96), PptxInches(3.72), PptxInches(5.1), PptxInches(0.25), "IMPLICATION", 9, PPTX_ORANGE, bold=True)
        text_box(slide, PptxInches(6.96), PptxInches(4.08), PptxInches(5.1), PptxInches(1.24), _clip(insight["implication"] or "À interpréter selon le contexte métier.", 340), 14, PPTX_INK, valign=MSO_ANCHOR.MIDDLE)

    for index, encoded in enumerate(images_b64 or [], 1):
        try:
            raw = base64.b64decode(encoded)
            image_w, image_h = _image_dimensions(raw, PptxInches(11.2), PptxInches(4.95))
        except Exception as exc:
            logger.warning("Visualisation %s illisible, diapositive ignorée : %s", index, exc)
            continue
        slide = base_slide("Visualisation", f"Figure {index}")
        left = int((width - image_w) / 2)
        top = PptxInches(1.83) + int((PptxInches(4.95) - image_h) / 2)
        slide.shapes.add_picture(io.BytesIO(raw), left, top, width=int(image_w), height=int(image_h))

    recommendations = report["recommandations"]
    for page_index in range(0, len(recommendations), 3):
        chunk = recommendations[page_index:page_index + 3]
        suffix = "" if page_index == 0 else " — suite"
        slide = base_slide("Décision", f"Plan d'action{suffix}")
        y = 1.85
        for recommendation in chunk:
            rectangle(slide, PptxInches(0.85), PptxInches(y), PptxInches(11.55), PptxInches(1.36), PPTX_PAPER, PPTX_LINE)
            rectangle(slide, PptxInches(1.1), PptxInches(y + 0.22), PptxInches(1.34), PptxInches(0.34), PPTX_OLIVE)
            text_box(slide, PptxInches(1.13), PptxInches(y + 0.26), PptxInches(1.28), PptxInches(0.16), recommendation["priorite"].upper(), 7.5, PPTX_WHITE, bold=True, align=PP_ALIGN.CENTER)
            text_box(slide, PptxInches(2.72), PptxInches(y + 0.18), PptxInches(8.95), PptxInches(0.42), _clip(recommendation["action"], 150), 16, PPTX_INK, bold=True)
            text_box(slide, PptxInches(2.72), PptxInches(y + 0.69), PptxInches(8.95), PptxInches(0.42), _clip(recommendation["justification"], 220), 11.5, PPTX_MUTED)
            y += 1.58

    if report["limites"] or report["prochaines_etapes"]:
        slide = base_slide("Suite", "Limites et prochaines étapes")
        for x, heading, items, fill, color in (
            (0.85, "LIMITES", report["limites"], PPTX_PAPER, PPTX_ORANGE),
            (6.72, "PROCHAINES ÉTAPES", report["prochaines_etapes"], PPTX_OLIVE_LIGHT, PPTX_OLIVE),
        ):
            rectangle(slide, PptxInches(x), PptxInches(1.9), PptxInches(5.55), PptxInches(4.55), fill, PPTX_LINE)
            text_box(slide, PptxInches(x + 0.3), PptxInches(2.22), PptxInches(4.95), PptxInches(0.3), heading, 9, color, bold=True)
            y = 2.83
            for index, item in enumerate(items[:4], 1):
                text_box(slide, PptxInches(x + 0.34), PptxInches(y), PptxInches(0.38), PptxInches(0.25), f"{index:02d}", 10, color, bold=True)
                text_box(slide, PptxInches(x + 0.82), PptxInches(y - 0.03), PptxInches(4.28), PptxInches(0.62), _clip(item, 155), 12.5, PPTX_INK)
                y += 0.82

    buffer = io.BytesIO()
    presentation.save(buffer)
    buffer.seek(0)
    return buffer.read()
